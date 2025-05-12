import streamlit as st
import pandas as pd
import requests
from bs4 import BeautifulSoup
import pypdf
import pptx
import openai
import io
import json
import re
import time
from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from openpyxl.utils import get_column_letter
from urllib.parse import urlparse, urlunparse

# --- Configuration ---
st.set_page_config(page_title="AI Marketing Content Generator", layout="wide")

# Load OpenAI API Key
try:
    openai.api_key = st.secrets["OPENAI_API_KEY"]
    client = openai.OpenAI(api_key=openai.api_key)
    MODEL = "gpt-4o-mini" # Use the specific model name
except Exception as e:
    st.error(f"Error loading OpenAI API key from secrets.toml: {e}")
    st.stop()

# --- Helper Functions ---

def add_http(url):
    """Adds http:// if scheme is missing from URL."""
    if not url:
        return url
    parsed = urlparse(url)
    if not parsed.scheme:
        # Check if it looks like a domain name (e.g., contains a dot)
        if '.' in parsed.path.split('/')[0]:
             # Basic check, assumes 'domain.com/path' if no scheme
            return f"https://{url}" # Default to https
        else:
            # If it doesn't look like a domain (e.g., just 'example'), return as is or handle differently
            return url # Or raise an error, or return None
    return urlunparse(parsed)

@st.cache_data(ttl=3600) # Cache for 1 hour
def extract_text_from_url(url):
    """Extracts text content from a URL."""
    try:
        headers = {'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/58.0.3029.110 Safari/537.3'}
        response = requests.get(url, headers=headers, timeout=20)
        response.raise_for_status()  # Raise an exception for bad status codes
        soup = BeautifulSoup(response.content, 'lxml') # Use lxml for better parsing

        # Remove script and style elements
        for script_or_style in soup(["script", "style"]):
            script_or_style.decompose()

        # Get text, strip whitespace, and join lines
        text = ' '.join(soup.stripped_strings)
        return text
    except requests.exceptions.RequestException as e:
        st.warning(f"Could not fetch URL {url}: {e}")
        return None
    except Exception as e:
        st.warning(f"Error parsing URL {url}: {e}")
        return None

@st.cache_data(ttl=3600)
def extract_text_from_pdf(file):
    """Extracts text content from an uploaded PDF file."""
    try:
        reader = pypdf.PdfReader(file)
        text = ""
        for page in reader.pages:
            page_text = page.extract_text()
            if page_text:
                text += page_text + "\n"
        return text
    except Exception as e:
        st.warning(f"Error reading PDF file {file.name}: {e}")
        return None

@st.cache_data(ttl=3600)
def extract_text_from_ppt(file):
    """Extracts text content from an uploaded PPT file."""
    try:
        prs = pptx.Presentation(file)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    except Exception as e:
        st.warning(f"Error reading PPT file {file.name}: {e}")
        return None

def summarize_text(text, max_chars=3000, source_name=""):
    """Summarizes text using OpenAI API."""
    if not text or not text.strip():
        return f"No content extracted from {source_name}." if source_name else "No content provided."

    prompt = f"""
    Please summarize the following text extracted from '{source_name}'.
    Focus on the company's core offerings, value proposition, target audience, and key marketing messages.
    The summary should be concise and informative, suitable for generating marketing ad copy.
    Keep the summary under {max_chars} characters.

    Text to summarize:
    ---
    {text[:40000]}
    ---

    Summary:
    """ # Limit input text slightly to avoid hitting token limits easily

    try:
        response = client.chat.completions.create(
            model=MODEL,
            messages=[
                {"role": "system", "content": "You are a helpful assistant skilled in summarizing business context for marketing purposes."},
                {"role": "user", "content": prompt}
            ],
            max_tokens=1000, # Adjust based on expected summary length
            temperature=0.5,
        )
        summary = response.choices[0].message.content.strip()
        return summary
    except Exception as e:
        st.error(f"Error summarizing text from {source_name}: {e}")
        return f"Error during summarization for {source_name}."

def generate_ad_content(prompt):
    """Generates ad content using OpenAI API, expecting JSON output."""
    try:
        response = client.chat.completions.create(
            model=MODEL,
            messages=[
                {"role": "system", "content": "You are an expert marketing copywriter. Generate ad content based on the provided context and instructions. Output *only* valid JSON."},
                {"role": "user", "content": prompt}
            ],
            max_tokens=4000, # Allow ample tokens for JSON generation
            temperature=0.7,
            response_format={"type": "json_object"}, # Enforce JSON output if model supports it
        )
        content = response.choices[0].message.content.strip()
        # Basic validation if JSON is returned
        try:
            json_content = json.loads(content)
            return json_content
        except json.JSONDecodeError:
            st.error(f"AI did not return valid JSON. Raw response:\n```\n{content}\n```")
            # Attempt to extract JSON from potential markdown code blocks
            match = re.search(r'```json\s*([\s\S]*?)\s*```', content, re.IGNORECASE)
            if match:
                try:
                    json_content = json.loads(match.group(1))
                    st.warning("Extracted JSON from markdown block.")
                    return json_content
                except json.JSONDecodeError:
                    st.error("Failed to parse extracted JSON.")
                    return None
            return None # Indicate failure
    except Exception as e:
        st.error(f"Error generating AI content: {e}")
        return None

def create_styled_excel(data_dict, company_name, lead_objective_str):
    """Creates a styled Excel file from the generated ad data."""
    filename_safe_company_name = re.sub(r'[\\/*?:"<>|]', "", company_name) # Sanitize filename
    filename_safe_objective = lead_objective_str.replace(" ", "_")
    filename = f"{filename_safe_company_name}_{filename_safe_objective}.xlsx"

    output = io.BytesIO()
    with pd.ExcelWriter(output, engine='openpyxl') as writer:
        # Define styles
        header_font = Font(color="FFFFFF", bold=True)
        header_fill = PatternFill(start_color="000000", end_color="000000", fill_type="solid")
        header_alignment = Alignment(horizontal='center', vertical='center')
        content_alignment = Alignment(vertical='middle', wrap_text=True)
        thin_border_side = Side(border_style="thin", color="000000")
        cell_border = Border(left=thin_border_side, right=thin_border_side, top=thin_border_side, bottom=thin_border_side)

        # --- Write Sheets ---
        for sheet_name, df in data_dict.items():
            if df is not None and not df.empty:
                df.to_excel(writer, sheet_name=sheet_name, index=False)
                worksheet = writer.sheets[sheet_name]

                # Apply styles
                for row in worksheet.iter_rows():
                    for cell in row:
                        cell.alignment = content_alignment
                        cell.border = cell_border

                # Style Header Row
                for cell in worksheet[1]:
                    cell.font = header_font
                    cell.fill = header_fill
                    cell.alignment = header_alignment

                # Adjust column widths (basic approach)
                for col_idx, column in enumerate(df.columns):
                    max_length = 0
                    column_letter = get_column_letter(col_idx + 1)

                    # Check header length
                    if df[column].name:
                       max_length = max(max_length, len(str(df[column].name)))

                    # Check cell content lengths
                    for cell in df[column]:
                        try:
                            if cell is not None:
                                # Add a buffer for wrapped text lines
                                cell_len = max(len(line) for line in str(cell).split('\n'))
                                max_length = max(max_length, cell_len)
                        except:
                            pass # Ignore errors in length calculation

                    # Set width (add padding) - max width around 70-80 to prevent huge columns
                    adjusted_width = min((max_length + 5) * 1.2, 70)
                    worksheet.column_dimensions[column_letter].width = adjusted_width

            else:
                # Create an empty sheet if no data was generated
                workbook = writer.book
                if sheet_name not in workbook.sheetnames:
                     workbook.create_sheet(sheet_name)
                worksheet = workbook[sheet_name]
                worksheet['A1'] = f"No data generated for {sheet_name}"
                # Apply basic styling if needed

    output.seek(0)
    return output, filename

# --- Prompt Design Functions ---

def create_email_prompt(summary, count, objective_link):
    return f"""
    Based on the following company context summary:
    ---
    {summary}
    ---
    Generate {count} distinct email ad variations for a 'Demand Capture' objective.
    Each variation should include:
    1.  Ad Name: A unique identifier (up to 250 chars), e.g., "Email_DemandCapture_Variant_1_Topic".
    2.  Objective: Always "Demand Capture".
    3.  Headline: A compelling email headline (often similar to Subject Line but can be used internally).
    4.  Subject Line: An engaging subject line to maximize open rates.
    5.  Body: 2-3 paragraphs of persuasive copy. Embed the objective link '{objective_link}' naturally within the text (e.g., using markdown link format like [click here]({objective_link}) or similar phrasing). Focus on the value proposition and encourage the reader to take the next step.
    6.  CTA: A concise call-to-action phrase related to the body's main message (e.g., "Book Your Demo", "Schedule a Meeting").

    Output the result as a JSON object with a single key "emails", which is an array of objects, each object representing an email variation with keys "Ad Name", "Objective", "Headline", "Subject Line", "Body", "CTA".

    Example JSON structure:
    {{
      "emails": [
        {{
          "Ad Name": "Email_DemandCapture_Variant_1_FeatureX",
          "Objective": "Demand Capture",
          "Headline": "Unlock Feature X Benefits",
          "Subject Line": "Ready to boost your results with Feature X?",
          "Body": "Paragraph 1 introducing the problem and solution.\\n\\nParagraph 2 highlighting benefits and value. Learn more and [book your demo]({objective_link}) today!\\n\\nParagraph 3 reinforcing the call to action.",
          "CTA": "Book Your Demo Now"
        }},
        // ... more email objects
      ]
    }}
    """

def create_linkedin_prompt(summary, count, objective, destination_link, cta_button_text):
    return f"""
    Based on the following company context summary:
    ---
    {summary}
    ---
    Generate {count} distinct LinkedIn ad variations for the objective: '{objective}'.
    Each variation should include:
    1.  Ad Name: A unique identifier (up to 250 chars), e.g., "LinkedIn_{objective.replace(' ','')}_Variant_1_Topic".
    2.  Objective: "{objective}".
    3.  Introductory Text: 300-400 characters. The first 150 characters must contain a strong hook. Include 1-2 relevant emojis naturally.
    4.  Image Copy: Suggest concise text (1-2 short phrases or bullet points) that could overlay an ad image, reinforcing the main message.
    5.  Headline: Around 70 characters, compelling and clear.
    6.  Destination: The link "{destination_link}".
    7.  CTA Button: "{cta_button_text}".

    Output the result as a JSON object with a single key "linkedin_ads", which is an array of objects, each representing a LinkedIn ad variation with keys "Ad Name", "Objective", "Introductory Text", "Image Copy", "Headline", "Destination", "CTA Button".

    Example JSON structure:
    {{
      "linkedin_ads": [
        {{
          "Ad Name": "LinkedIn_{objective.replace(' ','')}_Variant_1_BenefitY",
          "Objective": "{objective}",
          "Introductory Text": "Struggling with [Problem]? 🤔 Discover how [Company/Product] helps you achieve [Benefit]. Our solution offers [Value Prop]. Learn more! 👇 #Marketing #AdTech",
          "Image Copy": "Achieve [Benefit] Faster | [Key Feature]",
          "Headline": "Stop Guessing, Start Growing: Achieve [Benefit] Today",
          "Destination": "{destination_link}",
          "CTA Button": "{cta_button_text}"
        }},
        // ... more LinkedIn ad objects
      ]
    }}
    """

def create_facebook_prompt(summary, count, objective, destination_link, cta_button_text):
     return f"""
    Based on the following company context summary:
    ---
    {summary}
    ---
    Generate {count} distinct Facebook ad variations for the objective: '{objective}'.
    Each variation should include:
    1.  Ad Name: A unique identifier (up to 250 chars), e.g., "Facebook_{objective.replace(' ','')}_Variant_1_Topic".
    2.  Objective: "{objective}".
    3.  Primary Text: 300-400 characters. The first 125 characters must contain a strong hook. Include 1-2 relevant emojis naturally.
    4.  Image Copy: Suggest concise text (1-2 short phrases or bullet points) that could overlay an ad image, reinforcing the main message.
    5.  Headline: Around 27 characters, punchy and attention-grabbing.
    6.  Link Description: Around 27 characters, providing context for the link.
    7.  Destination: The link "{destination_link}".
    8.  CTA Button: "{cta_button_text}".

    Output the result as a JSON object with a single key "facebook_ads", which is an array of objects, each representing a Facebook ad variation with keys "Ad Name", "Objective", "Primary Text", "Image Copy", "Headline", "Link Description", "Destination", "CTA Button".

    Example JSON structure:
    {{
      "facebook_ads": [
        {{
          "Ad Name": "Facebook_{objective.replace(' ','')}_Variant_1_BenefitZ",
          "Objective": "{objective}",
          "Primary Text": "Tired of [Pain Point]? 😫 See how [Company/Product] makes [Task] easy! Get [Result] without the hassle. Click below to find out more! ✨",
          "Image Copy": "[Benefit] Made Simple | Try Us Free",
          "Headline": "Unlock [Benefit] Now!",
          "Link Description": "Click here for details!",
          "Destination": "{destination_link}",
          "CTA Button": "{cta_button_text}"
        }},
        // ... more Facebook ad objects
      ]
    }}
    """

def create_google_search_prompt(summary):
    return f"""
    Based on the following company context summary:
    ---
    {summary}
    ---
    Generate content for Google Search Responsive Search Ads (RSAs). Provide:
    1.  Headlines: Exactly 15 unique headlines, each around 30 characters maximum. Focus on keywords, benefits, and calls to action.
    2.  Descriptions: Exactly 4 unique descriptions, each around 90 characters maximum. Elaborate on value propositions and encourage clicks.

    Output the result as a JSON object with two keys: "headlines" (an array of 15 strings) and "descriptions" (an array of 4 strings).

    Example JSON structure:
    {{
      "headlines": [
        "Headline 1 (Max 30 Chars)",
        "Headline 2 - Benefit",
        "Headline 3 - Keyword",
        "Headline 4 - CTA",
        "Headline 5 - Location?",
        "Headline 6 - Offer",
        "Headline 7 - Urgency",
        "Headline 8 - Brand Name",
        "Headline 9 - Feature",
        "Headline 10 - Question?",
        "Headline 11 - Social Proof",
        "Headline 12 - Unique Value",
        "Headline 13 - Simplicity",
        "Headline 14 - Result",
        "Headline 15 - Alternative CTA"
      ],
      "descriptions": [
        "Description 1: Elaborate on key benefit and include a call to action. (Max 90 Chars)",
        "Description 2: Highlight a different feature or value prop. Mention target audience maybe.",
        "Description 3: Focus on problem/solution and unique selling points. Add credibility.",
        "Description 4: Combine benefits and provide a clear next step for the user. Urgency/Offer."
      ]
    }}
    """

def create_google_display_prompt(summary):
    return f"""
    Based on the following company context summary:
    ---
    {summary}
    ---
    Generate content for Google Display Responsive Display Ads. Provide:
    1.  Headlines: Exactly 5 unique headlines, each around 30 characters maximum. Focus on grabbing attention visually.
    2.  Descriptions: Exactly 5 unique descriptions, each around 90 characters maximum. Provide more context and encourage clicks.

    Output the result as a JSON object with two keys: "headlines" (an array of 5 strings) and "descriptions" (an array of 5 strings).

    Example JSON structure:
    {{
      "headlines": [
        "Display Headline 1 (Max 30)",
        "Benefit-Focused Headling",
        "Short & Punchy Headline",
        "Intriguing Question?",
        "Clear Call to Action"
      ],
      "descriptions": [
        "Display Description 1: Explain the offer or value proposition clearly. (Max 90 Chars)",
        "Display Description 2: Focus on a key benefit relevant to display audiences.",
        "Display Description 3: Highlight what makes the company/product unique.",
        "Display Description 4: Encourage exploration or learning more.",
        "Display Description 5: A concise summary with a call to action."
      ]
    }}
    """


# --- Streamlit App UI ---

st.title("🚀 AI Marketing Content Generator")
st.markdown("Extract context from your website and documents to generate tailored ad copy.")

# --- Inputs ---
st.header("1. Provide Company Context")
company_url = st.text_input("Client's Website URL (e.g., https://www.example.com)", "")
uploaded_context_files = st.file_uploader(
    "Upload Additional Context (PDF or PPT)",
    type=["pdf", "pptx"],
    accept_multiple_files=True
)
uploaded_lead_magnet = st.file_uploader(
    "Upload Downloadable Lead Magnet (PDF)",
    type=["pdf"]
)

st.header("2. Configure Ad Options")
lead_objective = st.selectbox(
    "Primary Lead Objective",
    ["Demo Booking", "Sales Meeting", "Lead Magnet Download"] # Added Lead Magnet option
)
learn_more_link = st.text_input("Link for 'Learn More' CTAs (e.g., product page, homepage)", "")
download_link_input = st.text_input("Link to Downloadable Material (if different from uploaded PDF)", "")
objective_link_input = st.text_input(f"Link for '{lead_objective}' Objective (e.g., Calendly, Hubspot meeting link, Thank You page)", "")

content_count = st.slider("Number of Ad Variations per Type/Objective", 1, 15, 10) # Reduced max to 10 for performance

st.header("3. Generate Content")
generate_button = st.button("✨ Generate Ad Content")

# --- Backend Flow ---
if generate_button:
    # Validate inputs
    if not company_url and not uploaded_context_files:
        st.error("Please provide a Website URL or upload at least one context file.")
        st.stop()
    if lead_objective == "Lead Magnet Download" and not uploaded_lead_magnet and not download_link_input:
         st.error("Please upload a Lead Magnet PDF or provide a Download Link for the 'Lead Magnet Download' objective.")
         st.stop()
    if lead_objective != "Lead Magnet Download" and not objective_link_input:
        st.error(f"Please provide the Link for the '{lead_objective}' objective.")
        st.stop()
    if not learn_more_link:
        st.warning("Consider adding a 'Learn More' link for Brand Awareness campaigns.")


    # Determine links based on objective
    download_link = download_link_input # Use provided link first
    if not download_link and uploaded_lead_magnet:
        # Placeholder - actual serving of uploaded file needs more infrastructure
        # For now, we'll just use the input link or leave it blank if neither is provided
        st.warning("Using uploaded lead magnet requires hosting. Provide 'Link to Downloadable Material' for generated ads.")
        # download_link = f"https://example.com/download/{uploaded_lead_magnet.name}" # Placeholder

    objective_link = objective_link_input

    # Derive company name from URL if possible
    company_name = "CompanyName"
    if company_url:
        try:
            parsed_url = urlparse(add_http(company_url))
            domain = parsed_url.netloc
            if domain.startswith('www.'):
                domain = domain[4:]
            company_name = domain.split('.')[0].capitalize() # Basic extraction
        except Exception:
            company_name = "CompanyName" # Fallback

    # Initialize progress
    progress_bar = st.progress(0)
    status_text = st.empty()
    total_steps = 6 # Extraction, Summarization, Email, LinkedIn, Facebook, Google, Excel

    extracted_texts = {}
    # 1. Extract Context
    status_text.text("Step 1/6: Extracting context...")
    if company_url:
        url_text = extract_text_from_url(add_http(company_url))
        if url_text:
            extracted_texts['website'] = url_text
            st.info(f"Extracted ~{len(url_text)} characters from {company_url}")
        else:
            st.warning(f"Could not extract text from {company_url}")

    if uploaded_context_files:
        for file in uploaded_context_files:
            file_content = None
            if file.type == "application/pdf":
                file_content = extract_text_from_pdf(file)
            elif file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
                file_content = extract_text_from_ppt(file)

            if file_content:
                extracted_texts[file.name] = file_content
                st.info(f"Extracted ~{len(file_content)} characters from {file.name}")
            else:
                 st.warning(f"Could not extract text from {file.name}")

    if not extracted_texts:
        st.error("No text could be extracted from the provided sources. Cannot proceed.")
        st.stop()

    progress_bar.progress(1/total_steps)

    # 2. Summarize Context
    status_text.text("Step 2/6: Summarizing context...")
    summaries = []
    with st.spinner("AI is summarizing the extracted content..."):
        for name, text in extracted_texts.items():
            # Adjust max_chars based on input length?
            max_chars = 3000 if len(text) > 5000 else 1800
            summary = summarize_text(text, max_chars, name)
            summaries.append(f"--- Context from {name} ---\n{summary}")
            time.sleep(1) # Small delay between API calls

    combined_summary = "\n\n".join(summaries)

    if not combined_summary or all("Error" in s for s in summaries):
         st.error("Failed to summarize the context. Cannot proceed.")
         st.stop()

    st.subheader("Combined Summary for Ad Generation:")
    st.text_area("Summary", combined_summary, height=200)
    progress_bar.progress(2/total_steps)

    # 3. Design Prompts & Generate Content
    generated_data = {}
    all_ads_data = {} # To store dataframes

    # --- Email ---
    status_text.text("Step 3/6: Generating Email Ads...")
    with st.spinner("Generating Email content..."):
        email_prompt = create_email_prompt(combined_summary, content_count, objective_link)
        email_json = generate_ad_content(email_prompt)
        if email_json and 'emails' in email_json:
            try:
                all_ads_data['Email'] = pd.DataFrame(email_json['emails'])
                st.success("✅ Email content generated.")
            except Exception as e:
                st.error(f"Error creating Email DataFrame: {e}")
                st.json(email_json) # Show raw JSON if parsing fails
                all_ads_data['Email'] = pd.DataFrame() # Empty DF
        else:
            st.warning("Could not generate Email content.")
            all_ads_data['Email'] = pd.DataFrame()
    progress_bar.progress(3/total_steps)
    time.sleep(1)

    # --- LinkedIn ---
    status_text.text("Step 4/6: Generating LinkedIn Ads...")
    linkedin_dfs = []
    linkedin_objectives = ["Brand Awareness", "Demand Gen", "Demand Capture"]
    with st.spinner("Generating LinkedIn content..."):
        for obj in linkedin_objectives:
            status_text.text(f"Step 4/6: Generating LinkedIn Ads ({obj})...")
            dest_link = learn_more_link
            cta = "Learn More"
            if obj == "Demand Gen":
                dest_link = download_link
                cta = "Download"
            elif obj == "Demand Capture":
                dest_link = objective_link
                cta = "Register" if lead_objective == "Demo Booking" else "Request Demo" # Or Book Now? Adjust as needed

            if not dest_link: # Fallback if a specific link is missing
                 dest_link = learn_more_link if learn_more_link else company_url

            linkedin_prompt = create_linkedin_prompt(combined_summary, content_count, obj, dest_link, cta)
            linkedin_json = generate_ad_content(linkedin_prompt)
            if linkedin_json and 'linkedin_ads' in linkedin_json:
                 try:
                    linkedin_dfs.append(pd.DataFrame(linkedin_json['linkedin_ads']))
                 except Exception as e:
                    st.error(f"Error creating LinkedIn DataFrame for {obj}: {e}")
                    st.json(linkedin_json)
            else:
                st.warning(f"Could not generate LinkedIn content for {obj}.")
            time.sleep(1) # Delay

        if linkedin_dfs:
            all_ads_data['LinkedIn'] = pd.concat(linkedin_dfs, ignore_index=True)
            st.success("✅ LinkedIn content generated.")
        else:
            all_ads_data['LinkedIn'] = pd.DataFrame()
    progress_bar.progress(4/total_steps)


    # --- Facebook ---
    status_text.text("Step 5/6: Generating Facebook Ads...")
    facebook_dfs = []
    facebook_objectives = ["Brand Awareness", "Demand Gen", "Demand Capture"]
    with st.spinner("Generating Facebook content..."):
        for obj in facebook_objectives:
            status_text.text(f"Step 5/6: Generating Facebook Ads ({obj})...")
            dest_link = learn_more_link
            cta = "Learn More"
            if obj == "Demand Gen":
                dest_link = download_link
                cta = "Download"
            elif obj == "Demand Capture":
                dest_link = objective_link
                cta = "Book Now" # Common FB CTA

            if not dest_link: # Fallback
                 dest_link = learn_more_link if learn_more_link else company_url

            facebook_prompt = create_facebook_prompt(combined_summary, content_count, obj, dest_link, cta)
            facebook_json = generate_ad_content(facebook_prompt)
            if facebook_json and 'facebook_ads' in facebook_json:
                try:
                    facebook_dfs.append(pd.DataFrame(facebook_json['facebook_ads']))
                except Exception as e:
                    st.error(f"Error creating Facebook DataFrame for {obj}: {e}")
                    st.json(facebook_json)
            else:
                st.warning(f"Could not generate Facebook content for {obj}.")
            time.sleep(1) # Delay

        if facebook_dfs:
            all_ads_data['FaceBook'] = pd.concat(facebook_dfs, ignore_index=True) # Note sheet name change
            st.success("✅ Facebook content generated.")
        else:
             all_ads_data['FaceBook'] = pd.DataFrame()
    progress_bar.progress(5/total_steps)


    # --- Google Search & Display ---
    status_text.text("Step 6/6: Generating Google Ads & Excel...")
    google_search_data = None
    google_display_data = None

    with st.spinner("Generating Google Search & Display content..."):
        # Google Search
        search_prompt = create_google_search_prompt(combined_summary)
        search_json = generate_ad_content(search_prompt)
        if search_json and 'headlines' in search_json and 'descriptions' in search_json:
            try:
                # Pad shorter list to make DataFrame creation easier
                max_len = max(len(search_json['headlines']), len(search_json['descriptions']))
                headlines = search_json['headlines'] + [None] * (max_len - len(search_json['headlines']))
                descriptions = search_json['descriptions'] + [None] * (max_len - len(search_json['descriptions']))
                google_search_data = pd.DataFrame({'Headline': headlines, 'Description': descriptions})
                all_ads_data['Google Search'] = google_search_data
                st.success("✅ Google Search content generated.")
            except Exception as e:
                st.error(f"Error creating Google Search DataFrame: {e}")
                st.json(search_json)
                all_ads_data['Google Search'] = pd.DataFrame()
        else:
            st.warning("Could not generate Google Search content.")
            all_ads_data['Google Search'] = pd.DataFrame()
        time.sleep(1)

        # Google Display
        display_prompt = create_google_display_prompt(combined_summary)
        display_json = generate_ad_content(display_prompt)
        if display_json and 'headlines' in display_json and 'descriptions' in display_json:
             try:
                max_len = max(len(display_json['headlines']), len(display_json['descriptions']))
                headlines = display_json['headlines'] + [None] * (max_len - len(display_json['headlines']))
                descriptions = display_json['descriptions'] + [None] * (max_len - len(display_json['descriptions']))
                google_display_data = pd.DataFrame({'Headline': headlines, 'Description': descriptions})
                all_ads_data['Google Display'] = google_display_data
                st.success("✅ Google Display content generated.")
             except Exception as e:
                st.error(f"Error creating Google Display DataFrame: {e}")
                st.json(display_json)
                all_ads_data['Google Display'] = pd.DataFrame()
        else:
            st.warning("Could not generate Google Display content.")
            all_ads_data['Google Display'] = pd.DataFrame()

    # 4. Parse and Format into XLSX
    status_text.text("Formatting Excel file...")
    try:
        excel_bytes, excel_filename = create_styled_excel(all_ads_data, company_name, lead_objective)
        progress_bar.progress(1.0)
        status_text.success("🎉 Content Generation Complete!")

        # 5. Enable Download
        st.download_button(
            label="⬇️ Download Ad Content (.xlsx)",
            data=excel_bytes,
            file_name=excel_filename,
            mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
        )

        # Optionally display generated dataframes in the app
        st.subheader("Generated Content Preview:")
        for name, df in all_ads_data.items():
            if df is not None and not df.empty:
                st.markdown(f"### {name}")
                st.dataframe(df)
            else:
                 st.markdown(f"### {name}")
                 st.write(f"No data generated for {name}.")


    except Exception as e:
        st.error(f"Error creating Excel file: {e}")
        progress_bar.progress(1.0)
        status_text.error("An error occurred during Excel file generation.")